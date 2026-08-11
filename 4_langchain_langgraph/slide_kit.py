"""Slide template for Voltway Research, our fictional fleet analyst brand."""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image  # add this import at the top of the file

NAVY = RGBColor(0x0A, 0x25, 0x40)
TEAL = RGBColor(0x00, 0xC2, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xB9, 0xC4, 0xD0)

LOGO = os.path.join(os.path.dirname(__file__), "sandbox", "skills", "fleet-slide", "logo.png")


def _textbox(slide, left, top, width, height, text, size, color, bold=False, spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    if spacing is not None:
        para.line_spacing = spacing
    return box


def build_slide(title: str, key_points: list[str], recommendation: str, outfile: str) -> None:
    """Build the one-slide Voltway recommendation deck and save it to outfile."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    slide.shapes.add_picture(LOGO, Inches(0.5), Inches(0.45), Inches(0.55), Inches(0.55))
    _textbox(slide, 1.2, 0.52, 6, 0.5, "VOLTWAY  RESEARCH", 15, MIST, bold=True)

    _textbox(slide, 0.5, 1.55, 12.3, 1.2, title, 34, WHITE, bold=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.53), Inches(2.55), Inches(2.2), Inches(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()

    for i, point in enumerate(key_points[:3]):
        y = 3.05 + i * 0.72
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y + 0.1), Inches(0.16), Inches(0.16))
        marker.fill.solid()
        marker.fill.fore_color.rgb = TEAL
        marker.line.fill.background()
        _textbox(slide, 0.95, y, 11.8, 0.6, point, 17, WHITE)

    _textbox(slide, 0.53, 5.55, 6, 0.4, "VOLTWAY RECOMMENDS", 12, TEAL, bold=True)
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.95), Inches(12.33), Inches(1.05))
    banner.adjustments[0] = 0.18
    banner.fill.solid()
    banner.fill.fore_color.rgb = TEAL
    banner.line.fill.background()
    frame = banner.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.35)
    frame.margin_right = Inches(0.35)
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = recommendation
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = NAVY

    prs.save(outfile)

def build_slide_new(
    title: str,
    key_points: list[str],
    recommendation: str,
    outfile: str,
    car_image: str | None = None,
) -> None:
    """Build the one-slide Voltway recommendation deck and save it to outfile."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    slide.shapes.add_picture(LOGO, Inches(0.5), Inches(0.45), Inches(0.55), Inches(0.55))
    _textbox(slide, 1.2, 0.52, 6, 0.5, "VOLTWAY  RESEARCH", 15, MIST, bold=True)

    # Content column width shrinks to make room for the image on the right
    content_width = 6.7 if car_image else 12.3

    _textbox(slide, 0.5, 1.55, content_width, 1.2, title, 34, WHITE, bold=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.53), Inches(2.55), Inches(2.2), Inches(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL
    rule.line.fill.background()

    for i, point in enumerate(key_points[:3]):
        y = 3.05 + i * 0.72
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(y + 0.1), Inches(0.16), Inches(0.16))
        marker.fill.solid()
        marker.fill.fore_color.rgb = TEAL
        marker.line.fill.background()
        _textbox(slide, 0.95, y, content_width - 0.4, 0.6, point, 17, WHITE)

    # Right-half car image, aspect-ratio preserved and centered in its box
    if car_image:
        box_left, box_top = 7.4, 1.55
        box_width, box_height = 5.4, 4.0  # available right-half area above the banner

        with Image.open(car_image) as img:
            img_w, img_h = img.size
        img_ratio = img_w / img_h
        box_ratio = box_width / box_height

        if img_ratio > box_ratio:
            # image is relatively wider than the box -> fit to width
            draw_width = box_width
            draw_height = draw_width / img_ratio
        else:
            # image is relatively taller than the box -> fit to height
            draw_height = box_height
            draw_width = draw_height * img_ratio

        left = box_left + (box_width - draw_width) / 2
        top = box_top + (box_height - draw_height) / 2

        slide.shapes.add_picture(car_image, Inches(left), Inches(top), Inches(draw_width), Inches(draw_height))

    _textbox(slide, 0.53, 5.55, 6, 0.4, "VOLTWAY RECOMMENDS", 12, TEAL, bold=True)
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.95), Inches(12.33), Inches(1.05))
    banner.adjustments[0] = 0.18
    banner.fill.solid()
    banner.fill.fore_color.rgb = TEAL
    banner.line.fill.background()
    frame = banner.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.35)
    frame.margin_right = Inches(0.35)
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = recommendation
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = NAVY

    prs.save(outfile)